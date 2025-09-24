# Copyright 2025 DataRobot, Inc.
#
# Licensed under the Apache License, Version 2.0 (the "License");
# You may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#   http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.
"""
Document loader for RAG-Ultra: efficiently extracts per-page text from PDFs, DOCX, PPTX, and TXT files.
Supports lazy loading, parallel processing, and robust fallback strategies for different file types.
"""

# TODO: Ask Brett: why not textract to support more file types?
import logging
import tempfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Callable, Dict, Tuple

import docx
import fitz  # PyMuPDF
import pptx
from fsspec import AbstractFileSystem

from ..persistent_fs.dr_file_system import get_file_system
from .constants import DEFAULT_MAX_WORKERS, SUPPORTED_FILE_TYPES, TEXT_FILE_TYPES
from .exceptions import (
    DocProcessorNoExtractorError,
    DocProcessorUnsupportedFileTypeError,
)

logger = logging.getLogger(__name__)


def convert_document_to_text(
    document_path: str,
    max_workers: int = DEFAULT_MAX_WORKERS,
    file_system: AbstractFileSystem | None = None,
) -> Dict[int, str]:
    """
    Extract per-page text from a document, auto-detecting file type.

    Args:
        document_path: Path to the document file.
        max_workers: Maximum number of worker threads for parallel processing.
        file_system: implementation of AbstractFileSystem for accessing to files, LocalFileSystem is default
    Returns:
        Dict mapping page numbers (1-indexed) to extracted text.
    Raises:
        ValueError: If document type is not supported.
        FileNotFoundError: If document file doesn't exist.
    """

    if not file_system:
        file_system = get_file_system()
    if not file_system.exists(document_path):
        raise FileNotFoundError(f"Document not found at {document_path}")

    path = Path(document_path)
    file_ext = path.suffix.lower().lstrip(".")

    if file_ext not in SUPPORTED_FILE_TYPES:
        raise DocProcessorUnsupportedFileTypeError(file_ext)
    if file_ext not in FILE_TYPES_TO_EXTRACTORS:
        raise DocProcessorNoExtractorError(file_ext)

    logger.info(f"Processing {file_ext} document: {document_path}")
    with tempfile.TemporaryDirectory() as tmpdirname:
        tmp_path = Path(tmpdirname) / path.name
        file_system.get(
            document_path, str(tmp_path)
        )  # copy file from persistent FS so we process locally

        return FILE_TYPES_TO_EXTRACTORS[file_ext](tmp_path, max_workers)


def _extract_pdf_page_fitz(path: Path, page_idx: int) -> Tuple[int, str]:
    """
    Helper for parallel PDF extraction using PyMuPDF.
    Returns (1-indexed page number, text).
    """
    try:
        with fitz.open(path) as doc:
            page = doc[page_idx]
            text = page.get_text()
            return (page_idx + 1, text)
    except Exception as e:
        logger.exception(f"Error extracting text from PDF page {page_idx + 1}: {e}")
        return (page_idx + 1, "")


def extract_text_from_pdf(
    path: Path, max_workers: int = DEFAULT_MAX_WORKERS
) -> Dict[int, str]:
    """
    Extract text from each page of a PDF using parallel processing.
    Prefers PyMuPDF for speed, falls back to PyPDF2 if needed.

    Args:
        path: Path to the PDF file.
        max_workers: Maximum number of worker threads.
    Returns:
        Dict mapping page numbers to page text.
    Raises:
        ImportError: If no PDF extraction library is available.
    """
    with fitz.open(path) as doc:
        page_count = len(doc)
    page_text = {}
    actual_workers = min(max_workers, max(1, page_count))
    with ThreadPoolExecutor(max_workers=actual_workers) as executor:
        future_to_page = {
            executor.submit(_extract_pdf_page_fitz, path, page_idx): page_idx
            for page_idx in range(page_count)
        }
        for future in as_completed(future_to_page):
            page_num, text = future.result()
            page_text[page_num] = text
    logger.info(f"Extracted text from {len(page_text)} PDF pages using PyMuPDF")
    return page_text


def extract_text_from_docx(
    path: Path, max_workers: int = DEFAULT_MAX_WORKERS
) -> Dict[int, str]:
    """
    Extract text from a DOCX file, splitting by page breaks.
    Each section between page breaks is treated as a "page".

    Args:
        path: Path to the Word document.
    Returns:
        Dict mapping simulated page numbers to text.
    Raises:
        ImportError: If python-docx is not installed.
    """
    # TODO: Add support for parallel processing if needed
    page_text = {}
    try:
        doc = docx.Document(str(path))
        current_page = 1
        current_text = ""
        for para in doc.paragraphs:
            if (
                "PAGE BREAK" in para.text.upper()
                or para.text.strip() == "\f"
                or (
                    hasattr(para, "style")
                    and para.style
                    and "page break" in str(para.style).lower()
                )
            ):
                if current_text.strip():
                    page_text[current_page] = current_text.strip()
                    current_page += 1
                    current_text = ""
            else:
                current_text += para.text + "\n"
        if current_text.strip():
            page_text[current_page] = current_text.strip()
        logger.info(f"Extracted {len(page_text)} pages from DOCX document")
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {e}")
        raise
    return page_text


def extract_text_from_pptx(
    path: Path, max_workers: int = DEFAULT_MAX_WORKERS
) -> Dict[int, str]:
    """
    Extract text from a PPTX file, treating each slide as a page.

    Args:
        path: Path to the PowerPoint presentation.
    Returns:
        Dict mapping slide numbers to slide text.
    Raises:
        ImportError: If python-pptx is not installed.
    """
    page_text = {}
    try:
        presentation = pptx.Presentation(str(path))
        for i, slide in enumerate(presentation.slides):
            text_list = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_list.append(shape.text)
            page_text[i + 1] = "\n".join(text_list)
        logger.info(f"Extracted text from {len(page_text)} slides")
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {e}")
        raise
    return page_text


def extract_text_from_txt(
    path: Path, max_workers: int = DEFAULT_MAX_WORKERS
) -> Dict[int, str]:
    """
    Extract text from a TXT file, splitting by page markers or length.

    Args:
        path: Path to the text file.
    Returns:
        Dict mapping page numbers to page text.
    Raises:
        Exception: If file cannot be read or split.
    """
    try:
        with open(path, "r", encoding="utf-8") as file:
            content = file.read()
        pages = split_text_into_pages(content)
        page_text = {
            i + 1: page.strip() for i, page in enumerate(pages) if page.strip()
        }
        return page_text
    except Exception as e:
        logger.error(f"Error extracting text from TXT: {e}")
        raise


def split_text_into_pages(content: str, max_chars_per_page: int = 3000) -> list[str]:
    """
    Split text into pages using common page markers or by paragraph length.
    Tries to preserve natural breaks and keep page sizes reasonable.

    Args:
        content: Text content to split.
        max_chars_per_page: Maximum characters per page.
    Returns:
        List of page content strings.
    """
    page_markers = ["\f", "----", "****", "======", "# Page", "===", "---", "***"]
    for marker in page_markers:
        if marker in content:
            pages = content.split(marker)
            logger.info(f"Split text file by marker: {marker}")
            return pages
    paragraphs = content.split("\n\n")
    pages = []
    current_page = ""
    for para in paragraphs:
        if len(current_page) + len(para) > max_chars_per_page and current_page:
            pages.append(current_page)
            current_page = para
        else:
            if current_page:
                current_page += "\n\n" + para
            else:
                current_page = para
    if current_page:
        pages.append(current_page)
    logger.info(f"Split text file into {len(pages)} pages by paragraph breaks")
    return pages


FILE_TYPES_TO_EXTRACTORS: Dict[str, Callable[[Path, int], Dict[int, str]]] = {
    "pdf": extract_text_from_pdf,
    "docx": extract_text_from_docx,
    "pptx": extract_text_from_pptx,
    **dict.fromkeys(list(TEXT_FILE_TYPES), extract_text_from_txt),
}
